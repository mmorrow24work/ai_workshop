#!/usr/bin/env python3
"""Build a PowerPoint slide deck for the AI Workshop Zero to Hero Guide.

Structure: Title · Agenda · 12 chapter sections (section divider + content slides)
           · YouTube channels · AI Projects · Closing   (~42 slides total)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = Path(__file__).parent / "AI_Workshop_Slides.pptx"

# ─── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x5C, 0x8A)
MED_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0x44, 0x72, 0xC4)
PALE_BLUE  = RGBColor(0xBE, 0xD3, 0xF0)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x2C, 0x3E, 0x50)
MID_GREY   = RGBColor(0x60, 0x60, 0x60)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF8)
GHOST_NUM  = RGBColor(0x25, 0x72, 0xAD)  # watermark chapter number
RED_WARN   = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x27, 0xAE, 0x60)

FONT = "Arial"

# ─── Slide dimensions: 16:9 widescreen ─────────────────────────────────────
W        = Inches(13.33)
H        = Inches(7.5)
HEADER_H = Inches(1.2)
BODY_TOP = Inches(1.4)
MARGIN   = Inches(0.55)
COL_GAP  = Inches(0.35)


# ─── Low-level primitives ───────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def box(slide, x, y, w, h, color, border=False):
    """Add a solid-filled rectangle with no border by default."""
    shp = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not border:
        shp.line.fill.background()
    return shp


def tf(slide, x, y, w, h):
    """Return a text frame for a new text box."""
    frame = slide.shapes.add_textbox(x, y, w, h).text_frame
    frame.word_wrap = True
    return frame


def run(para, text, size=18, bold=False, italic=False, color=DARK_GREY):
    r = para.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def text(slide, x, y, w, h, txt, size=18, bold=False, italic=False,
         color=DARK_GREY, align=PP_ALIGN.LEFT):
    """Single-paragraph text box."""
    frame = tf(slide, x, y, w, h)
    p = frame.paragraphs[0]
    p.alignment = align
    run(p, txt, size=size, bold=bold, italic=italic, color=color)
    return frame


def multiline(slide, x, y, w, h, lines, size=18, bold=False, italic=False,
              color=DARK_GREY, align=PP_ALIGN.LEFT, line_space=6):
    """Multi-line text box — one paragraph per item in lines list."""
    frame = tf(slide, x, y, w, h)
    first = True
    for line in lines:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.alignment = align
        p.space_before = Pt(line_space)
        run(p, line, size=size, bold=bold, italic=italic, color=color)
    return frame


# ─── Bullet helper ──────────────────────────────────────────────────────────

def bullets(slide, items, x=None, y=None, w=None, h=None, base_size=19):
    """
    items — list of (level, text_str).
    level 0 = main bullet  ▸
    level 1 = sub-bullet   –
    level 2 = detail       ·
    """
    x = MARGIN        if x is None else x
    y = BODY_TOP      if y is None else y
    w = W - MARGIN*2  if w is None else w
    h = H - y - Inches(0.25) if h is None else h

    CHARS = {0: "▸", 1: "–", 2: "·"}
    INDENT = {0: "", 1: "    ", 2: "        "}

    frame = tf(slide, x, y, w, h)
    first = True
    for level, item_text in items:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        pfx = f"{INDENT[level]}{CHARS.get(level,'▸')}  "
        p.space_before = Pt(9 if level == 0 else 3)
        p.space_after  = Pt(0)
        size  = base_size - level * 2
        color = DARK_GREY if level == 0 else MID_GREY
        run(p, pfx + item_text, size=size, color=color)
    return frame


# ─── Shared slide furniture ─────────────────────────────────────────────────

def header(slide, title, badge=None):
    """Dark-blue header bar with optional right-aligned chapter badge."""
    box(slide, 0, 0, W, HEADER_H, DARK_BLUE)
    box(slide, 0, HEADER_H - Inches(0.06), W, Inches(0.06), ORANGE)
    title_w = W - MARGIN*2 - (Inches(1.8) if badge else 0)
    text(slide, MARGIN, Inches(0.18), title_w, Inches(0.9),
         title, size=26, bold=True, color=WHITE)
    if badge:
        bx = W - Inches(1.75)
        box(slide, bx, Inches(0.22), Inches(1.55), Inches(0.76), MED_BLUE)
        text(slide, bx, Inches(0.22), Inches(1.55), Inches(0.76),
             badge, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def col_header(slide, label, x, y, w):
    """Coloured sub-heading label above a column of bullets."""
    text(slide, x, y, w, Inches(0.38), label, size=13, bold=True, color=MED_BLUE)
    box(slide, x, y + Inches(0.40), w, Inches(0.04), LIGHT_BLUE)


# ─── Two-column bullet layout ───────────────────────────────────────────────

def two_columns(slide, left, right, left_label=None, right_label=None,
                y=None, base_size=18):
    y    = BODY_TOP if y is None else y
    cw   = (W - MARGIN*2 - COL_GAP) / 2
    x_l  = MARGIN
    x_r  = MARGIN + cw + COL_GAP

    y_content = y
    if left_label:
        col_header(slide, left_label, x_l, y, cw)
        y_content = y + Inches(0.52)
    bullets(slide, left,  x=x_l, y=y_content, w=cw,
            h=H - y_content - Inches(0.25), base_size=base_size)

    y2_content = y
    if right_label:
        col_header(slide, right_label, x_r, y, cw)
        y2_content = y + Inches(0.52)
    bullets(slide, right, x=x_r, y=y2_content, w=cw,
            h=H - y2_content - Inches(0.25), base_size=base_size)


# ─── Table helper ───────────────────────────────────────────────────────────

def table(slide, headers, rows, x=None, y=None, w=None, col_widths=None):
    x = MARGIN   if x is None else x
    y = BODY_TOP if y is None else y
    w = W - MARGIN*2 if w is None else w
    h = Inches(0.5 * (len(rows) + 1) + 0.2)

    tbl = slide.shapes.add_table(len(rows)+1, len(headers), x, y, w, h).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    # Header row
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, hdr, size=13, bold=True, color=WHITE)

    # Data rows
    for r_idx, row_data in enumerate(rows):
        bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx+1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            run(p, str(cell_text), size=12, color=DARK_GREY)

    return tbl


# ─── Named slide types ──────────────────────────────────────────────────────

def title_slide(prs):
    slide = blank(prs)
    box(slide, 0, 0, W, H, DARK_BLUE)
    # Subtle right panel
    box(slide, W * 0.60, 0, W * 0.40, H, RGBColor(0x16, 0x4F, 0x7A))
    box(slide, W * 0.60 - Inches(0.05), 0, Inches(0.05), H, MED_BLUE)
    box(slide, 0, H - Inches(0.14), W, Inches(0.14), ORANGE)

    text(slide, MARGIN, Inches(1.8), Inches(7.8), Inches(2.0),
         "AI Workshop", size=80, bold=True, color=WHITE)
    text(slide, MARGIN, Inches(3.85), Inches(7.8), Inches(1.0),
         "Zero to Hero Guide", size=34, color=PALE_BLUE)
    text(slide, MARGIN, Inches(5.0), Inches(7.8), Inches(0.5),
         "June 2026  ·  mmorrow24work/ai_workshop",
         size=14, italic=True, color=PALE_BLUE)

    multiline(slide, W*0.63 + Inches(0.3), Inches(2.2), Inches(4.5), Inches(4.0),
              ["12 Chapters", "", "For IT Professionals", "& Curious Beginners"],
              size=22, color=PALE_BLUE, align=PP_ALIGN.CENTER, line_space=8)


def section_slide(prs, num, title, tagline=""):
    slide = blank(prs)
    box(slide, 0, 0, W, H, DARK_BLUE)
    box(slide, 0, H - Inches(0.14), W, Inches(0.14), ORANGE)
    box(slide, 0, Inches(2.1), Inches(0.14), Inches(3.4), ORANGE)

    # Ghost chapter number watermark
    text(slide, Inches(5.5), Inches(0.3), Inches(8.0), Inches(6.5),
         num, size=280, bold=True, color=GHOST_NUM, align=PP_ALIGN.RIGHT)

    text(slide, Inches(0.4), Inches(2.1), Inches(10.0), Inches(0.55),
         f"Chapter {num}", size=16, italic=True, color=PALE_BLUE)
    text(slide, Inches(0.4), Inches(2.7), Inches(9.5), Inches(2.2),
         title, size=50, bold=True, color=WHITE)
    if tagline:
        text(slide, Inches(0.4), Inches(4.9), Inches(9.5), Inches(0.8),
             tagline, size=22, italic=True, color=PALE_BLUE)


def agenda_slide(prs, chapters):
    slide = blank(prs)
    header(slide, "What We'll Cover Today")
    mid = len(chapters) // 2
    two_columns(
        slide,
        [(0, f"{n}  {t}") for n, t in chapters[:mid]],
        [(0, f"{n}  {t}") for n, t in chapters[mid:]],
        base_size=17, y=BODY_TOP + Inches(0.1)
    )


def closing_slide(prs):
    slide = blank(prs)
    box(slide, 0, 0, W, H, DARK_BLUE)
    box(slide, 0, H - Inches(0.14), W, Inches(0.14), ORANGE)

    text(slide, MARGIN, Inches(1.3), W - MARGIN*2, Inches(1.2),
         "Where to Start", size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(slide, Inches(2), Inches(2.55), W - Inches(4), Inches(0.06), MED_BLUE)

    steps = [
        "1.  Try claude.ai or ChatGPT for a real task — today",
        "2.  Run  ollama run llama3  to see your first local model",
        "3.  Subscribe to The Batch (deeplearning.ai) newsletter",
        "4.  Choose ONE tool and use it for real work this week",
        "5.  Read the Security chapter before touching client data",
    ]
    frame = tf(slide, Inches(2.0), Inches(2.9), W - Inches(4.2), Inches(3.6))
    first = True
    for step in steps:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.space_before = Pt(14)
        run(p, step, size=20, color=WHITE)

    text(slide, MARGIN, H - Inches(0.95), W - MARGIN*2, Inches(0.5),
         "github.com/mmorrow24work/ai_workshop",
         size=14, italic=True, color=PALE_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  CHAPTER SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def ch01_slides(prs):
    # ── Section divider ──────────────────────────────────────────────────
    section_slide(prs, "01", "AI Interfaces",
                  "Six ways to talk to an LLM — from browser to API")

    # ── Content 1: The six interfaces ───────────────────────────────────
    slide = blank(prs)
    header(slide, "Six Ways to Talk to an LLM", "Chapter 01")
    two_columns(slide,
        [(0, "Browser — ChatGPT, Claude.ai, Gemini"),
         (1, "No setup required; best for quick tasks"),
         (0, "Desktop App — Claude Desktop, ChatGPT App"),
         (1, "MCP support; integrates with local files"),
         (0, "IDE Plugin — Copilot, Cursor, Continue.dev"),
         (1, "AI inside your code editor with full context")],
        [(0, "CLI — Claude Code, Ollama, Aider"),
         (1, "Scriptable; works inside CI/CD pipelines"),
         (0, "API — Anthropic, OpenAI, Google APIs"),
         (1, "Build applications; pay per token"),
         (0, "MCP — Model Context Protocol"),
         (1, "Plugin system connecting AI to local tools")],
        base_size=17)

    # ── Content 2: Comparison ────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Interface Comparison", "Chapter 01")
    table(slide,
          ["Interface", "Skill Level", "Best For", "Data Risk"],
          [["Browser",    "None",        "Quick questions",      "Medium"],
           ["Desktop App","Low",         "Daily use + MCP",      "Medium"],
           ["IDE Plugin", "Low–Medium",  "Coding assistance",    "Medium–High"],
           ["CLI",        "Medium–High", "Automation / DevOps",  "Low (local option)"],
           ["API",        "High",        "Building apps",        "Configurable"],
           ["MCP",        "Medium",      "Rich integrations",    "Configurable"]],
          col_widths=[Inches(2.2), Inches(2.0), Inches(4.5), Inches(3.5)])


def ch02_slides(prs):
    section_slide(prs, "02", "LLM History",
                  "From the transformer paper to the AI era — 2017 to 2026")

    slide = blank(prs)
    header(slide, "Key Milestones in LLM History", "Chapter 02")
    bullets(slide, [
        (0, "2017  —  \"Attention Is All You Need\" — Google's Transformer paper"),
        (1, "Parallel processing on GPUs; foundation of every modern LLM"),
        (0, "2018–19  —  BERT (Google) · GPT-1 and GPT-2 (OpenAI)"),
        (1, "GPT-2 so capable at text generation, OpenAI initially withheld it"),
        (0, "2020  —  GPT-3 (175B parameters) — few-shot learning breakthrough"),
        (0, "Nov 2022  —  ChatGPT launches — 1 million users in 5 days"),
        (0, "2023  —  GPT-4 · Claude · Gemini · Llama 2 · Mistral 7B"),
        (0, "2024–26  —  Reasoning models · agents · 1M token contexts · local inference"),
    ])


def ch03_slides(prs):
    section_slide(prs, "03", "Frontier LLMs in 2026",
                  "Proprietary and open-weight models — how to choose")

    # ── Proprietary ─────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Proprietary Frontier Models", "Chapter 03")
    two_columns(slide,
        [(0, "Anthropic — Claude Family"),
         (1, "Opus 4.8 · Sonnet 4.6 · Haiku 4.5"),
         (1, "200 K token context window"),
         (1, "Strength: long docs, safety, coding"),
         (0, "OpenAI — GPT & o-Series"),
         (1, "GPT-4o (multimodal) · o3 (reasoning)"),
         (1, "128 K context; o3 thinks before answering"),
         (1, "Strength: ecosystem, real-time voice")],
        [(0, "Google DeepMind — Gemini"),
         (1, "Gemini 2.0 Ultra · Flash"),
         (1, "Up to 1 MILLION token context"),
         (1, "Strength: Google Workspace integration"),
         (0, "xAI — Grok"),
         (1, "Real-time web access; integrated into X"),
         (1, "Less enterprise adoption than top three")],
        base_size=17)

    # ── Open-weight + choosing ───────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Open-Weight Models & Choosing Wisely", "Chapter 03")
    two_columns(slide,
        [(0, "Open-Weight: Free to download and run locally"),
         (0, "Meta — Llama 4"),
         (1, "Competes with proprietary on many benchmarks"),
         (0, "Mistral AI — Mistral 7B / Large"),
         (1, "Lean, efficient; good EU data-residency option"),
         (0, "Also: Qwen 2.5 · Phi-4 · DeepSeek V3")],
        [(0, "Decision Guide"),
         (0, "Max capability needed?"),
         (1, "Claude Opus or GPT o3"),
         (0, "Speed + cost balance?"),
         (1, "Claude Sonnet or GPT-4o"),
         (0, "Data must stay local?"),
         (1, "Ollama + Llama 4 or Mistral"),
         (0, "Google Workspace user?"),
         (1, "Gemini 2.0")],
        base_size=17)


def ch04_slides(prs):
    section_slide(prs, "04", "Offline LLM Hosting",
                  "Run AI locally — your data never leaves your machine")

    # ── Ollama ──────────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Ollama: Simplest Path to Local AI", "Chapter 04")
    two_columns(slide,
        [(0, "Why run locally?"),
         (1, "Data stays on your machine — no cloud risk"),
         (1, "No per-token API costs"),
         (1, "Works fully offline"),
         (0, "Install: curl -fsSL https://ollama.com/install.sh | sh"),
         (0, "ollama pull llama3"),
         (0, "ollama run llama3"),
         (0, "ollama run llama3 \"Explain REST APIs\"")],
        [(0, "Hardware guide"),
         (0, "7B model   → 8 GB RAM"),
         (0, "13B model  → 16 GB RAM"),
         (0, "34B model  → 32 GB RAM"),
         (0, "70B model  → 64 GB RAM"),
         (1, "Quantised models cut RAM needs by ~50 %"),
         (0, "Start with llama3 (8B) — fast and capable"),
         (0, "REST API at localhost:11434 (OpenAI-compatible)")],
        base_size=17)

    # ── Alternatives ─────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Local AI Alternatives", "Chapter 04")
    table(slide,
          ["Tool", "Interface", "Best For", "GPU Required?"],
          [["LM Studio",  "Desktop GUI", "Non-technical users; model browser",     "Optional"],
           ["Jan.ai",     "Desktop GUI", "Open-source; MCP support",               "Optional"],
           ["GPT4All",    "Desktop GUI", "CPU-only; older hardware",                "No"],
           ["LocalAI",    "API only",    "Developers; drop-in OpenAI replacement", "Optional"],
           ["Ollama",     "CLI + API",   "Developers; scriptable pipelines",       "Optional"]],
          col_widths=[Inches(1.8), Inches(2.0), Inches(5.5), Inches(2.3)])


def ch05_slides(prs):
    section_slide(prs, "05", "LLM Tokens",
                  "The currency of AI — what they are and why they matter")

    # ── What tokens are ──────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "What Is a Token?", "Chapter 05")
    two_columns(slide,
        [(0, "~4 characters  or  ¾ of an English word"),
         (0, "\"Hello, how are you?\"  →  ~5 tokens"),
         (0, "\"tokenisation\"  →  ~3 tokens"),
         (0, "Non-English text tokenises less efficiently"),
         (0, "Context Window = max tokens in + out"),
         (0, "Claude Opus/Sonnet/Haiku  —  200 K tokens"),
         (0, "GPT-4o  —  128 K tokens"),
         (0, "Gemini 2.0 Ultra  —  1 M tokens"),
         (1, "≈ an entire novel in a single conversation")],
        [(0, "Cost = input tokens + output tokens"),
         (0, "Claude Sonnet: ~$3 in / $15 out per 1M"),
         (0, "Claude Haiku:  ~$0.25 in / $1.25 out per 1M"),
         (0, "GPT-4o:        ~$5 in / $15 out per 1M"),
         (0, "GPT-4o-mini:   ~$0.15 in / $0.60 out per 1M"),
         (0, "Prompt caching reduces cost 80–90 %"),
         (1, "For repeated identical context/system prompts")],
        base_size=17)

    # ── Chunking ─────────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Managing Large Documents", "Chapter 05")
    bullets(slide, [
        (0, "When a document exceeds the context window, chunk it"),
        (0, "Fixed chunking — equal-sized pieces (e.g. 1 000 tokens with 200-token overlap)"),
        (0, "Semantic chunking — split at paragraphs, sections, or headings"),
        (0, "Hierarchical summarisation — summarise each chunk, then summarise summaries"),
        (0, "RAG — index the document and retrieve only relevant chunks at query time"),
        (1, "See Chapter 12 (Bonus Topics) for a full RAG walkthrough"),
        (0, "Rule of thumb: if it fits in context, include it — chunking adds complexity"),
    ])


def ch06_slides(prs):
    section_slide(prs, "06", "Subscription Packages",
                  "Claude Pro · ChatGPT Plus · Gemini Advanced — ~$20/month each")

    # ── Comparison ───────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "What $20/Month Gets You", "Chapter 06")
    table(slide,
          ["Feature", "Claude Pro", "ChatGPT Plus", "Gemini Advanced"],
          [["Best model",       "Opus 4.8",       "GPT-4o / o3",       "Gemini 2.0 Ultra"],
           ["Context window",   "200 K tokens",   "128 K tokens",      "1 M tokens"],
           ["Usage limit",      "~5 hrs heavy use","80 msgs / 3 hrs",  "Generous, undocumented"],
           ["Image generation", "✗",              "DALL-E 3",          "Imagen 3"],
           ["Web browsing",     "Limited",         "✓",                 "✓"],
           ["MCP support",      "✓  (Desktop)",   "Limited",           "✗"],
           ["Data training",    "✗  (Pro tier)",  "Opt-out available", "Opt-out available"]],
          col_widths=[Inches(2.8), Inches(2.7), Inches(2.7), Inches(2.7)])

    # ── API migration ─────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "When to Move from Subscription to API", "Chapter 06")
    two_columns(slide,
        [(0, "Move to API when…"),
         (0, "Building an application or automation"),
         (1, "Subscriptions are for human users only"),
         (0, "Need higher volume or custom rate limits"),
         (0, "Need fine-grained control"),
         (1, "temperature · max_tokens · system prompts"),
         (0, "Need programmatic integration with code"),
         (0, "Need per-token cost tracking")],
        [(0, "Getting started"),
         (0, "Anthropic: console.anthropic.com"),
         (1, "~$5 minimum credit to start"),
         (0, "OpenAI: platform.openai.com"),
         (0, "Google: aistudio.google.com"),
         (0, "Free tiers exist on all three APIs"),
         (0, "Never use free tiers with client data"),
         (1, "Free data may be used for model training")],
        base_size=17)


def ch07_slides(prs):
    section_slide(prs, "07", "AI Harness Options",
                  "Beyond the chat interface — tools built around the model")

    # ── The four main tools ───────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Developer Tools for AI-Assisted Work", "Chapter 07")
    two_columns(slide,
        [(0, "Claude Code  (CLI / Agent)"),
         (1, "npm install -g @anthropic-ai/claude-code"),
         (1, "Reads and writes your actual project files"),
         (1, "Git integration; runs tests; creates PRs"),
         (1, "Best: agentic multi-step coding tasks"),
         (0, "Cursor  (IDE — VS Code fork)"),
         (1, "AI-first editor; Composer for multi-file changes"),
         (1, "Tab-complete + Chat + Agent modes"),
         (1, "Best: developers wanting the smoothest UI")],
        [(0, "GitHub Copilot  (IDE Plugin)"),
         (1, "Works in VS Code, JetBrains, Neovim"),
         (1, "Inline completion + Copilot Chat"),
         (1, "Best: existing GitHub users; team licences"),
         (0, "Continue.dev  (Open-Source Plugin)"),
         (1, "VS Code + JetBrains; connect ANY model"),
         (1, "Works with local Ollama models (privacy!)"),
         (1, "Best: privacy-sensitive projects; cost control")],
        base_size=17)

    # ── Comparison ───────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Harness Comparison", "Chapter 07")
    table(slide,
          ["Tool", "Type", "Models", "Local Models?", "Price"],
          [["Claude Code",  "CLI / Agent",  "Claude only",  "No",  "API cost"],
           ["Cursor",       "IDE",          "Multi-model",  "No",  "Free / $20/mo"],
           ["Copilot",      "IDE Plugin",   "Multi-model",  "No",  "$10/mo"],
           ["Continue.dev", "IDE Plugin",   "Any model",    "Yes", "Free (BYOK)"],
           ["open-claw",    "Web UI",       "Claude only",  "No",  "API cost"],
           ["hermes-agent", "Framework",    "Flexible",     "Configurable", "Open-source"]],
          col_widths=[Inches(2.0), Inches(1.8), Inches(2.0), Inches(2.2), Inches(2.6)])


def ch08_slides(prs):
    section_slide(prs, "08", "Security Best Practices",
                  "Protecting client data when using AI tools")

    # ── Fundamental rule ──────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "The Fundamental Rule", "Chapter 08")
    # Warning box
    box(slide, MARGIN, BODY_TOP, W - MARGIN*2, Inches(1.4), RGBColor(0xFD, 0xED, 0xEC))
    box(slide, MARGIN, BODY_TOP, Inches(0.1), Inches(1.4), RED_WARN)
    text(slide, MARGIN + Inches(0.25), BODY_TOP + Inches(0.2),
         W - MARGIN*2 - Inches(0.4), Inches(1.0),
         "Never paste client data, PII, or confidential information into a public AI "
         "service without explicit organisational approval and a Data Processing Agreement.",
         size=18, bold=True, color=RED_WARN)
    bullets(slide, [
        (0, "Treat it like emailing a client spreadsheet to your personal Gmail"),
        (0, "Sensitive data types to protect:"),
        (1, "PII — names, email, phone, national IDs, dates of birth"),
        (1, "Financial — account numbers, transaction records, salaries"),
        (1, "Health data — diagnoses, medications, any medical information"),
        (1, "Client confidential — strategy, unreleased products, legal documents"),
        (1, "Credentials — passwords, API keys, tokens, private keys"),
    ], y=BODY_TOP + Inches(1.6))

    # ── API keys + prompt injection ───────────────────────────────────────
    slide = blank(prs)
    header(slide, "API Key Security & Prompt Injection", "Chapter 08")
    two_columns(slide,
        [(0, "API Key Security"),
         (0, "Store keys in environment variables"),
         (0, "Never commit keys to git repositories"),
         (0, "Use a secrets manager (Vault, AWS Secrets)"),
         (0, "Set spending limits on API accounts"),
         (0, "Use separate keys per project"),
         (0, "Rotate keys periodically")],
        [(0, "Prompt Injection Risk"),
         (1, "Malicious user input manipulates model behaviour"),
         (0, "Validate & sanitise all user-supplied input"),
         (0, "Keep system prompt separate from user content"),
         (0, "Never give AI admin access to critical systems"),
         (0, "For high-risk actions: require human confirmation"),
         (0, "Log all prompts and responses for audit trail")],
        base_size=17)

    # ── Checklist ─────────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Pre-Use Security Checklist", "Chapter 08")
    bullets(slide, [
        (0, "☐  Is this data covered by an NDA, privacy law, or client agreement?"),
        (0, "☐  Has this AI tool been approved by your organisation's security policy?"),
        (0, "☐  Is there a Data Processing Agreement with the AI provider?"),
        (0, "☐  Are you using an enterprise tier (data NOT used for model training)?"),
        (0, "☐  Is your API key in an env variable or secrets manager — not in code?"),
        (0, "☐  Have you considered a local model (Ollama) instead of a cloud service?"),
        (0, "☐  Is there a human reviewer in the loop for high-stakes outputs?"),
    ])


def ch09_slides(prs):
    section_slide(prs, "09", "AI Ethics",
                  "Building and using AI responsibly")

    slide = blank(prs)
    header(slide, "Key Ethical Dimensions", "Chapter 09")
    two_columns(slide,
        [(0, "Bias & Fairness"),
         (1, "Models learn patterns from biased training data"),
         (1, "Test across demographic groups before deployment"),
         (1, "Example: hiring model penalising postcodes ≈ race"),
         (0, "Hallucinations"),
         (1, "LLMs generate plausible but false information"),
         (1, "Always verify facts from authoritative sources"),
         (1, "Never rely on AI for legal citations without checking"),
         (0, "Consent & Transparency"),
         (1, "Users have a right to know they're talking to AI"),
         (1, "EU AI Act mandates disclosure in many contexts")],
        [(0, "Copyright & IP"),
         (1, "Training data includes copyrighted content"),
         (1, "Ownership of AI-generated content still contested"),
         (0, "Accessibility"),
         (1, "AI can include or exclude — design with care"),
         (1, "Consider users with disabilities, slow connections"),
         (0, "Environmental Impact"),
         (1, "Training a large model ≈ several transatlantic flights"),
         (1, "Prefer smaller or local models when frontier not needed"),
         (0, "Autonomy & Human Oversight"),
         (1, "Agentic AI needs meaningful human oversight"),
         (1, "High-stakes decisions must have a human in the loop")],
        base_size=16)


def ch10_slides(prs):
    section_slide(prs, "10", "AI Governance",
                  "Regulations, frameworks, and compliance — what IT teams must know")

    # ── Key frameworks ────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Key Regulatory Frameworks", "Chapter 10")
    two_columns(slide,
        [(0, "EU AI Act  (in force Aug 2024)"),
         (1, "World's first comprehensive AI law"),
         (1, "Applies to any AI used in the EU, anywhere built"),
         (0, "Risk tiers:"),
         (1, "Unacceptable  — Banned (social scoring, mass surveillance)"),
         (1, "High risk     — Strict requirements (hiring, credit, healthcare)"),
         (1, "Limited risk  — Transparency obligations (chatbots, deepfakes)"),
         (1, "Minimal risk  — No specific requirements"),
         (0, "GDPR  (existing, applies to AI on personal data)"),
         (1, "Art. 22: Right not to be subject to purely automated decisions")],
        [(0, "NIST AI RMF  (US; voluntary, growing in procurement)"),
         (1, "GOVERN · MAP · MEASURE · MANAGE"),
         (0, "ISO/IEC 42001  (published 2023; certifiable)"),
         (1, "Like ISO 27001 but for AI management systems"),
         (1, "Certification becoming a differentiator"),
         (0, "Australia — Privacy Act & AI Safety Standard"),
         (1, "APPs apply to AI handling personal data"),
         (1, "Voluntary AI Safety Standard (2024) issued"),
         (0, "Sector overlays: TGA (health) · ASIC (finance)"),
         (1, "Check sector-specific requirements for your domain")],
        base_size=16)

    # ── Compliance checklist ──────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Minimal Responsible AI Policy", "Chapter 10")
    two_columns(slide,
        [(0, "Your policy should cover:"),
         (0, "Approved use cases"),
         (0, "Prohibited use cases"),
         (0, "Data classification — what data meets which tools"),
         (0, "Approved tools list"),
         (0, "Human oversight requirements"),
         (0, "Incident response for AI failures"),
         (0, "Training requirements for staff"),
         (0, "Review cadence (recommend: quarterly)")],
        [(0, "Before deploying an AI system:"),
         (0, "☐  Identify risk level under EU AI Act"),
         (0, "☐  Conduct Data Protection Impact Assessment"),
         (0, "☐  Document purpose, inputs, and outputs"),
         (0, "☐  Establish human oversight for high-risk decisions"),
         (0, "☐  Create an audit trail for AI decisions"),
         (0, "☐  Brief all staff on acceptable use"),
         (0, "☐  Set a review date for the governance assessment")],
        base_size=17)


def ch11_slides(prs):
    section_slide(prs, "11", "Staying Current",
                  "AI moves fast — here's a sustainable system to keep up")

    slide = blank(prs)
    header(slide, "Best Sources & A Practical Reading System", "Chapter 11")
    two_columns(slide,
        [(0, "Newsletters (email)"),
         (0, "The Batch — deeplearning.ai"),
         (1, "Weekly; curated; excellent for IT professionals"),
         (0, "TLDR AI"),
         (1, "Daily; scannable in 5 minutes"),
         (0, "Ben's Bites"),
         (1, "Daily; strong on new tools and products"),
         (0, "Import AI — Jack Clark (Anthropic co-founder)"),
         (1, "Weekly; deep on research papers"),
         (0, "Communities"),
         (0, "r/LocalLlama · r/MachineLearning · HuggingFace Discord")],
        [(0, "A sustainable weekly system"),
         (0, "Daily (5 min)"),
         (1, "Scan TLDR AI headlines with morning coffee"),
         (0, "Weekly (30 min)"),
         (1, "Read The Batch digest"),
         (0, "Monthly (1–2 hrs)"),
         (1, "Deep-dive on ONE topic relevant to current work"),
         (0, "Quarterly"),
         (1, "Review tools — are there better options now?"),
         (0, "Most important habit:"),
         (1, "Apply what you learn. Reading without doing"),
         (1, "builds no real understanding.")],
        base_size=17)


def ch12_slides(prs):
    section_slide(prs, "12", "Bonus Topics",
                  "Prompt Engineering · RAG · AI Agents")

    # ── Prompt Engineering ────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Prompt Engineering: Get More From Every Query", "Chapter 12")
    two_columns(slide,
        [(0, "Be specific — format, length, audience"),
         (1, "\"Write a 3-paragraph summary for a Year 10 student\""),
         (0, "Assign a role via the system prompt"),
         (1, "\"You are a senior Python developer. Be concise.\""),
         (0, "Few-shot examples"),
         (1, "Show the pattern 2–3 times before the real task"),
         (0, "Chain-of-thought"),
         (1, "\"Think step by step before answering\"")],
        [(0, "Request structured output"),
         (1, "\"Return JSON with: summary, confidence (0–100), sources\""),
         (0, "Set explicit constraints"),
         (1, "\"Under 150 words. No brand names. If unsure, say so.\""),
         (0, "System prompt structure for apps:"),
         (1, "Role → Context → Instructions → Output format → Constraints"),
         (0, "Resource: docs.anthropic.com/prompt-engineering")],
        base_size=17)

    # ── RAG ───────────────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "RAG: Grounding AI in Your Own Documents", "Chapter 12")
    two_columns(slide,
        [(0, "The problem"),
         (1, "LLMs don't know your internal documents"),
         (1, "Training data has a cutoff date"),
         (0, "RAG solution: Index → Retrieve → Answer"),
         (0, "Step 1 — Indexing (done once)"),
         (1, "Split docs into chunks → embed → store in vector DB"),
         (0, "Step 2 — Querying (every request)"),
         (1, "Embed the question → find similar chunks"),
         (1, "Include chunks in prompt → LLM answers with context")],
        [(0, "Popular vector databases"),
         (0, "ChromaDB — local, simple, great to start"),
         (0, "Pinecone — managed cloud, production-ready"),
         (0, "pgvector — PostgreSQL extension"),
         (0, "Qdrant — high-performance open-source"),
         (0, "Key use cases"),
         (0, "Company knowledge base / internal docs"),
         (0, "Customer support with product manuals"),
         (0, "Policy Q&A grounded in real documents"),
         (0, "Frameworks: LangChain · LlamaIndex")],
        base_size=17)

    # ── AI Agents ─────────────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "AI Agents: From Answering to Acting", "Chapter 12")
    two_columns(slide,
        [(0, "Traditional LLM (single-turn)"),
         (1, "You ask  →  It answers  →  Done"),
         (0, "AI Agent (multi-turn, tool-using)"),
         (1, "Goal  →  Plan  →  Act  →  Observe  →  Iterate  →  Report"),
         (0, "Agents use tools"),
         (1, "Read/write files · run terminal commands"),
         (1, "Search the web · call APIs · query databases"),
         (0, "Example: Claude Code as an agent"),
         (1, "Plans which files to edit for a task"),
         (1, "Writes code · runs tests · commits — autonomously")],
        [(0, "Real-world agent examples"),
         (0, "Claude Code — coding + DevOps automation"),
         (0, "OpenHands — full software engineering agent"),
         (0, "AutoGPT — early autonomous task pioneer"),
         (0, "Multi-agent systems"),
         (1, "Orchestrator → Research Agent + Analysis Agent + Writer"),
         (0, "Key safety considerations"),
         (1, "Human confirmation for destructive actions"),
         (1, "Limit tool access to what's needed (least privilege)"),
         (1, "Set maximum iteration limits"),
         (1, "Log everything for auditability")],
        base_size=17)


# ─── Resource slides ─────────────────────────────────────────────────────────

def youtube_slide(prs):
    slide = blank(prs)
    header(slide, "Top 10 Must-Watch YouTube Channels")
    two_columns(slide,
        [(0, "Andrej Karpathy"),
         (1, "Neural Networks: Zero to Hero — build GPT from scratch"),
         (0, "3Blue1Brown"),
         (1, "Visual maths; best intro to neural networks anywhere"),
         (0, "Two Minute Papers"),
         (1, "3–5 min paper summaries — stay current effortlessly"),
         (0, "Yannic Kilcher"),
         (1, "Deep-dives on landmark AI research papers"),
         (0, "AI Explained"),
         (1, "Accessible model comparisons and capability breakdowns")],
        [(0, "Matt Wolfe"),
         (1, "Weekly new tools roundup — beginner-friendly"),
         (0, "Fireship"),
         (1, "Fast, funny developer perspective; 100-seconds format"),
         (0, "Sam Witteveen"),
         (1, "Hands-on LangChain, RAG, and agent tutorials"),
         (0, "Sentdex"),
         (1, "Python and ML fundamentals; friendly long-form"),
         (0, "Google DeepMind"),
         (1, "Research explainers and AlphaFold-style documentaries")],
        base_size=17)


def projects_slide(prs):
    slide = blank(prs)
    header(slide, "Top 10 AI Projects to Explore & Learn From")
    two_columns(slide,
        [(0, "Ollama  (github.com/ollama/ollama)"),
         (1, "Run LLMs locally in one command — start here"),
         (0, "Stable Diffusion WebUI  (AUTOMATIC1111)"),
         (1, "Open-source image generation; one of GitHub's most-starred"),
         (0, "Hugging Face Transformers"),
         (1, "The standard Python library for working with AI models"),
         (0, "LangChain"),
         (1, "Framework for building LLM apps, RAG, and agents"),
         (0, "LlamaIndex"),
         (1, "Specialised data ingestion and retrieval for LLMs")],
        [(0, "PrivateGPT"),
         (1, "100 % local document Q&A — nothing leaves your machine"),
         (0, "OpenHands (ex-OpenDevin)"),
         (1, "Open-source AI software engineering agent"),
         (0, "AnythingLLM"),
         (1, "Polished self-hosted AI stack with user management"),
         (0, "Perplexica"),
         (1, "Open-source Perplexity alternative — web-grounded answers"),
         (0, "Anthropic Cookbook"),
         (1, "Official code examples for every Claude API pattern")],
        base_size=17)


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build_pptx():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    chapters = [
        ("01", "AI Interfaces"),
        ("02", "LLM History"),
        ("03", "Frontier LLMs in 2026"),
        ("04", "Offline LLM Hosting"),
        ("05", "LLM Tokens"),
        ("06", "Subscription Packages"),
        ("07", "AI Harness Options"),
        ("08", "Security Best Practices"),
        ("09", "AI Ethics"),
        ("10", "AI Governance"),
        ("11", "Staying Current"),
        ("12", "Bonus Topics"),
    ]

    print("Building slides …")
    title_slide(prs);             print("  ✓ Title")
    agenda_slide(prs, chapters);  print("  ✓ Agenda")

    ch01_slides(prs);  print("  ✓ Ch 01 — AI Interfaces")
    ch02_slides(prs);  print("  ✓ Ch 02 — LLM History")
    ch03_slides(prs);  print("  ✓ Ch 03 — Frontier LLMs")
    ch04_slides(prs);  print("  ✓ Ch 04 — Offline Hosting")
    ch05_slides(prs);  print("  ✓ Ch 05 — Tokens")
    ch06_slides(prs);  print("  ✓ Ch 06 — Packages")
    ch07_slides(prs);  print("  ✓ Ch 07 — Harness Options")
    ch08_slides(prs);  print("  ✓ Ch 08 — Security")
    ch09_slides(prs);  print("  ✓ Ch 09 — Ethics")
    ch10_slides(prs);  print("  ✓ Ch 10 — Governance")
    ch11_slides(prs);  print("  ✓ Ch 11 — Staying Current")
    ch12_slides(prs);  print("  ✓ Ch 12 — Bonus Topics")

    youtube_slide(prs);  print("  ✓ YouTube Channels")
    projects_slide(prs); print("  ✓ AI Projects")
    closing_slide(prs);  print("  ✓ Closing")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_pptx()
